"""
PowerPoint Operations Helper Scripts

Tool functions for the powerpoint_operations skill.
Provides CRUD operations on PowerPoint presentations.
"""

from pathlib import Path
from typing import Optional, Dict, Any, List


def create_presentation(file_path: str, title: str = "Presentation", author: str = "") -> Dict[str, Any]:
    """
    Create a new PowerPoint presentation with a title slide.

    Args:
        file_path: Path where the .pptx file will be saved
        title: Presentation title
        author: Optional author name

    Returns:
        Dict with status and file info
    """
    from pptx import Presentation
    from pptx.util import Inches

    path = Path(file_path)
    path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.core_properties.title = title
    if author:
        prs.core_properties.author = author

    # Add title slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    prs.save(str(path))

    return {
        "status": "created",
        "file_path": str(path.resolve()),
        "title": title,
        "slide_count": 1,
        "size_bytes": path.stat().st_size,
    }


def add_slide(file_path: str, layout_index: int = 1, title: str = "", content: str = "") -> Dict[str, Any]:
    """
    Add a new slide to an existing presentation.

    Args:
        file_path: Path to the .pptx file
        layout_index: Slide layout (0=Title, 1=Title+Content, 2=Two Content, 5=Blank)
        title: Slide title text
        content: Slide body text (plain text, one line per bullet)

    Returns:
        Dict with status and slide info
    """
    from pptx import Presentation

    path = Path(file_path)
    if not path.exists():
        return {"error": f"File not found: {file_path}"}

    prs = Presentation(str(path))

    if layout_index >= len(prs.slide_layouts):
        return {"error": f"Layout index {layout_index} out of range. Max: {len(prs.slide_layouts) - 1}"}

    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)

    if title and slide.shapes.title:
        slide.shapes.title.text = title

    if content:
        # Find the body placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Body placeholder
                tf = shape.text_frame
                tf.text = ""
                for i, line in enumerate(content.split("\n")):
                    if i == 0:
                        tf.text = line
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                break

    prs.save(str(path))

    return {
        "status": "added",
        "slide_number": len(prs.slides),
        "layout_index": layout_index,
        "title": title,
    }


def read_slides(file_path: str) -> Dict[str, Any]:
    """
    Read all slides and extract their text content.

    Args:
        file_path: Path to the .pptx file

    Returns:
        Dict with slide list and text content
    """
    from pptx import Presentation

    path = Path(file_path)
    if not path.exists():
        return {"error": f"File not found: {file_path}"}

    prs = Presentation(str(path))
    slides_data = []

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        slides_data.append({
            "slide_number": i + 1,
            "text": "\n".join(texts),
            "shape_count": len(slide.shapes),
        })

    return {
        "file_path": str(path.resolve()),
        "slide_count": len(slides_data),
        "title": prs.core_properties.title or "",
        "author": prs.core_properties.author or "",
        "slides": slides_data,
    }


def get_presentation_info(file_path: str) -> Dict[str, Any]:
    """
    Get metadata about a PowerPoint presentation.

    Args:
        file_path: Path to the .pptx file

    Returns:
        Dict with metadata and structure info
    """
    from pptx import Presentation

    path = Path(file_path)
    if not path.exists():
        return {"error": f"File not found: {file_path}"}

    prs = Presentation(str(path))

    return {
        "file_path": str(path.resolve()),
        "title": prs.core_properties.title or "",
        "author": prs.core_properties.author or "",
        "subject": prs.core_properties.subject or "",
        "slide_count": len(prs.slides),
        "layout_count": len(prs.slide_layouts),
        "size_bytes": path.stat().st_size,
    }


def delete_slide(file_path: str, slide_number: int) -> Dict[str, Any]:
    """
    Delete a slide by its 1-based number.

    Args:
        file_path: Path to the .pptx file
        slide_number: 1-based slide number to delete

    Returns:
        Dict with status
    """
    from pptx import Presentation
    from pptx.oxml.ns import qn

    path = Path(file_path)
    if not path.exists():
        return {"error": f"File not found: {file_path}"}

    prs = Presentation(str(path))
    if slide_number < 1 or slide_number > len(prs.slides):
        return {"error": f"Slide {slide_number} out of range. Total slides: {len(prs.slides)}"}

    slide = prs.slides[slide_number - 1]
    sp = slide._element
    sp.getparent().remove(sp)

    prs.save(str(path))

    return {
        "status": "deleted",
        "deleted_slide": slide_number,
        "remaining_slides": len(prs.slides),
    }


def add_text_box(file_path: str, slide_number: int, text: str, left: float = 1.0, top: float = 1.0, width: float = 5.0, height: float = 1.0, font_size: int = 18) -> Dict[str, Any]:
    """
    Add a text box to a specific slide.

    Args:
        file_path: Path to the .pptx file
        slide_number: 1-based slide number
        text: Text content
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        font_size: Font size in points

    Returns:
        Dict with status
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    path = Path(file_path)
    if not path.exists():
        return {"error": f"File not found: {file_path}"}

    prs = Presentation(str(path))
    if slide_number < 1 or slide_number > len(prs.slides):
        return {"error": f"Slide {slide_number} out of range."}

    slide = prs.slides[slide_number - 1]
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.text = text

    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)

    prs.save(str(path))

    return {"status": "added", "slide_number": slide_number, "text": text[:50]}


def set_slide_background(file_path: str, slide_number: int, color_hex: str) -> Dict[str, Any]:
    """
    Set the background color of a slide.

    Args:
        file_path: Path to the .pptx file
        slide_number: 1-based slide number
        color_hex: Background color in hex (e.g., "FF5733" without #)

    Returns:
        Dict with status
    """
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from lxml import etree

    path = Path(file_path)
    if not path.exists():
        return {"error": f"File not found: {file_path}"}

    prs = Presentation(str(path))
    if slide_number < 1 or slide_number > len(prs.slides):
        return {"error": f"Slide {slide_number} out of range."}

    slide = prs.slides[slide_number - 1]
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color_hex)

    prs.save(str(path))

    return {"status": "background_set", "slide_number": slide_number, "color": color_hex}


def duplicate_slide(file_path: str, slide_number: int) -> Dict[str, Any]:
    """
    Duplicate a slide by its 1-based number.

    Args:
        file_path: Path to the .pptx file
        slide_number: 1-based slide number to duplicate

    Returns:
        Dict with status
    """
    from pptx import Presentation
    from copy import deepcopy

    path = Path(file_path)
    if not path.exists():
        return {"error": f"File not found: {file_path}"}

    prs = Presentation(str(path))
    if slide_number < 1 or slide_number > len(prs.slides):
        return {"error": f"Slide {slide_number} out of range."}

    source_slide = prs.slides[slide_number - 1]
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    # Copy shapes
    for shape in source_slide.shapes:
        el = deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    prs.save(str(path))

    return {
        "status": "duplicated",
        "source_slide": slide_number,
        "new_slide_number": len(prs.slides),
    }
