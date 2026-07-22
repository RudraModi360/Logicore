"""
PPTX auto-fixes.

Handles fixing common PowerPoint issues:
- Fix shape alignment (move back within boundaries)
- Add title to first slide
- Remove empty slides
"""

from __future__ import annotations

import logging
import os
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from logicore.verification.auto_fix import AutoFixEngine

logger = logging.getLogger(__name__)


def register_fixes(engine: "AutoFixEngine") -> None:
    """Register PPTX fix handlers with the auto-fix engine."""
    engine.register_fix("alignment", "extends beyond", _fix_shape_alignment)
    engine.register_fix("content", "no title text", _fix_add_title)
    engine.register_fix("content", "empty", _fix_remove_empty_slides)


def _fix_shape_alignment(artifact_path: str, issue) -> bool:
    """Fix shapes that extend beyond slide boundaries."""
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        return False

    try:
        prs = Presentation(artifact_path)
        modified = False

        for slide in prs.slides:
            try:
                slide_width = slide.slide_layout.slide_width
                slide_height = slide.slide_layout.slide_height
            except Exception:
                continue

            for shape in slide.shapes:
                try:
                    # Fix left boundary.
                    if shape.left is not None and shape.left < 0:
                        shape.left = Emu(0)
                        modified = True

                    # Fix right boundary.
                    if shape.left is not None and shape.width is not None:
                        right_edge = shape.left + shape.width
                        if right_edge > slide_width:
                            shape.left = slide_width - shape.width
                            modified = True

                    # Fix top boundary.
                    if shape.top is not None and shape.top < 0:
                        shape.top = Emu(0)
                        modified = True

                    # Fix bottom boundary.
                    if shape.top is not None and shape.height is not None:
                        bottom_edge = shape.top + shape.height
                        if bottom_edge > slide_height:
                            shape.top = slide_height - shape.height
                            modified = True

                except Exception:
                    continue

        if modified:
            prs.save(artifact_path)

        return modified

    except Exception as exc:
        logger.debug(f"Failed to fix shape alignment: {exc}")
        return False


def _fix_add_title(artifact_path: str, issue) -> bool:
    """Add a title text box to the first slide if missing."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        return False

    try:
        prs = Presentation(artifact_path)

        if len(prs.slides) == 0:
            return False

        first_slide = prs.slides[0]

        # Check if first slide already has title-like text.
        for shape in first_slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                return False  # Already has title.

        # Add a title text box.
        left = Inches(1)
        top = Inches(0.5)
        width = Inches(8)
        height = Inches(1.5)

        txBox = first_slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Presentation Title"

        # Style the title.
        paragraph = tf.paragraphs[0]
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 0, 0)

        prs.save(artifact_path)
        return True

    except Exception as exc:
        logger.debug(f"Failed to add title: {exc}")
        return False


def _fix_remove_empty_slides(artifact_path: str, issue) -> bool:
    """Remove slides that have no content."""
    try:
        from pptx import Presentation
    except ImportError:
        return False

    try:
        prs = Presentation(artifact_path)

        if len(prs.slides) <= 1:
            return False  # Don't remove if only one slide.

        slides_to_remove = []

        for i, slide in enumerate(prs.slides):
            # Check if slide has any meaningful content.
            has_content = False
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    has_content = True
                    break
                elif shape.shape_type is not None:
                    # Non-text shape exists.
                    has_content = True
                    break

            if not has_content:
                slides_to_remove.append(i)

        if not slides_to_remove:
            return False

        # Remove slides (from end to start).
        for i in reversed(slides_to_remove):
            slide_id = prs.slides._sldIdLst[i]
            prs.slides._sldIdLst.remove(slide_id)

        prs.save(artifact_path)
        return True

    except Exception as exc:
        logger.debug(f"Failed to remove empty slides: {exc}")
        return False
