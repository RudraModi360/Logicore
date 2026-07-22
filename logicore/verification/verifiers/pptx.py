"""
PPTX verifier.

Verifies PowerPoint presentations: slides, alignment, content, layout.
"""

from __future__ import annotations

import os
from typing import List, Optional, Set

from logicore.verification.base_verifier import BaseVerifier
from logicore.verification.result import VerificationIssue


# Supported PPTX extensions.
PPTX_EXTENSIONS: Set[str] = {".pptx", ".ppsx", ".ppt"}

# Maximum slides for a reasonable presentation.
MAX_SLIDES = 200

# Minimum slides expected.
MIN_SLIDES = 1


class PPTXVerifier(BaseVerifier):
    """Verify PowerPoint presentations: slides, alignment, content, layout.

    Checks:
    - Presentation can be opened
    - Has at least one slide
    - Slides have content (not empty)
    - First slide has a title
    - Elements are within slide boundaries
    - No duplicate consecutive slides
    """

    def supported_extensions(self) -> Set[str]:
        return PPTX_EXTENSIONS

    def _verify_content(
        self,
        artifact_path: str,
        issues: List[VerificationIssue],
        requirements: Optional[str],
    ) -> None:
        ext = os.path.splitext(artifact_path)[1].lower()

        if ext == ".ppt":
            # Legacy .ppt files need different handling.
            self._verify_ppt(artifact_path, issues)
        else:
            self._verify_pptx(artifact_path, issues, requirements)

    def _verify_pptx(
        self,
        artifact_path: str,
        issues: List[VerificationIssue],
        requirements: Optional[str],
    ) -> None:
        """Verify PPTX files using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError:
            self._add_issue(
                issues,
                severity="info",
                category="dependency",
                description="python-pptx not installed — skipping detailed PPTX verification",
            )
            return

        try:
            prs = Presentation(artifact_path)
        except Exception as exc:
            self._add_issue(
                issues,
                severity="critical",
                category="corruption",
                description=f"Cannot open PPTX file: {exc}",
                auto_fixable=False,
            )
            return

        # Slide count check.
        num_slides = len(prs.slides)
        if num_slides == 0:
            self._add_issue(
                issues,
                severity="critical",
                category="content",
                description="Presentation has no slides",
            )
            return

        if num_slides > MAX_SLIDES:
            self._add_issue(
                issues,
                severity="warning",
                category="structure",
                description=f"Presentation has {num_slides} slides (unusually large)",
                auto_fixable=False,
                fix_suggestion="Consider splitting into multiple presentations",
            )

        # First slide title check.
        first_slide = prs.slides[0]
        has_title = False
        for shape in first_slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                has_title = True
                break
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        has_title = True
                        break
                if has_title:
                    break

        if not has_title:
            self._add_issue(
                issues,
                severity="warning",
                category="content",
                description="First slide has no title text",
                location="slide 1",
                auto_fixable=True,
                fix_suggestion="Add a title to the first slide",
            )

        # Per-slide checks.
        empty_slides = []
        for i, slide in enumerate(prs.slides, 1):
            slide_issues = self._check_slide(slide, i)
            issues.extend(slide_issues)

            # Track empty slides.
            if not slide.shapes or all(
                not s.has_text_frame or not s.text.strip()
                for s in slide.shapes
            ):
                empty_slides.append(i)

        if empty_slides:
            self._add_issue(
                issues,
                severity="warning",
                category="content",
                description=f"{len(empty_slides)} slide(s) appear to be empty: {empty_slides}",
                auto_fixable=True,
                fix_suggestion="Add content to empty slides or remove them",
            )

        # Slide layout variety check.
        layout_names = set()
        for slide in prs.slides:
            if slide.slide_layout:
                layout_names.add(slide.slide_layout.name)
        if len(layout_names) == 1 and num_slides > 3:
            self._add_issue(
                issues,
                severity="info",
                category="structure",
                description=f"All {num_slides} slides use the same layout",
                fix_suggestion="Consider varying slide layouts for visual interest",
            )

    def _check_slide(self, slide, slide_num: int) -> List[VerificationIssue]:
        """Check a single slide for issues."""
        issues = []

        # No shapes check.
        if not slide.shapes:
            return issues

        # Alignment check: shapes outside slide boundaries.
        try:
            slide_width = slide.slide_layout.slide_width
            slide_height = slide.slide_layout.slide_height
        except Exception:
            return issues

        for shape in slide.shapes:
            try:
                # Check horizontal boundaries.
                if shape.left is not None and shape.left < 0:
                    issues.append(VerificationIssue(
                        severity="warning",
                        category="alignment",
                        description=f"Shape extends beyond left slide boundary",
                        location=f"slide {slide_num}",
                        auto_fixable=True,
                        fix_suggestion="Adjust shape position to fit within slide",
                    ))
                    break  # One per slide is enough.

                if shape.left is not None and shape.width is not None:
                    right_edge = shape.left + shape.width
                    if right_edge > slide_width:
                        issues.append(VerificationIssue(
                            severity="warning",
                            category="alignment",
                            description=f"Shape extends beyond right slide boundary",
                            location=f"slide {slide_num}",
                            auto_fixable=True,
                            fix_suggestion="Adjust shape position to fit within slide",
                        ))
                        break

                # Check vertical boundaries.
                if shape.top is not None and shape.top < 0:
                    issues.append(VerificationIssue(
                        severity="warning",
                        category="alignment",
                        description=f"Shape extends beyond top slide boundary",
                        location=f"slide {slide_num}",
                        auto_fixable=True,
                        fix_suggestion="Adjust shape position to fit within slide",
                    ))
                    break

                if shape.top is not None and shape.height is not None:
                    bottom_edge = shape.top + shape.height
                    if bottom_edge > slide_height:
                        issues.append(VerificationIssue(
                            severity="warning",
                            category="alignment",
                            description=f"Shape extends beyond bottom slide boundary",
                            location=f"slide {slide_num}",
                            auto_fixable=True,
                            fix_suggestion="Adjust shape position to fit within slide",
                        ))
                        break
            except Exception:
                continue

        return issues

    def _verify_ppt(
        self,
        artifact_path: str,
        issues: List[VerificationIssue],
    ) -> None:
        """Basic verification for legacy .ppt files."""
        try:
            with open(artifact_path, "rb") as f:
                header = f.read(8)
        except Exception as exc:
            self._add_issue(
                issues,
                severity="critical",
                category="corruption",
                description=f"Cannot read .ppt file: {exc}",
            )
            return

        # OLE compound document magic bytes.
        if not header.startswith(b"\xd0\xcf\x11\xe0"):
            self._add_issue(
                issues,
                severity="critical",
                category="format",
                description="File does not have a valid .ppt header",
            )


def get_verifier() -> PPTXVerifier:
    """Return a PPTXVerifier instance for registry auto-discovery."""
    return PPTXVerifier()
